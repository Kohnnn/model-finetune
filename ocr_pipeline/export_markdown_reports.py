from __future__ import annotations

import argparse
import hashlib
import json
import logging
import re
import sys
from pathlib import Path
from typing import Any

try:
    from .process_pdfs import (
        build_document_metadata,
        discover_files,
        extract_docx_pages,
        extract_pdf_pages,
        extract_pptx_pages,
        normalize_extensions,
        normalize_text,
        select_files,
        strip_head_boilerplate,
        strip_vcsc_disclaimers,
        trim_tail_sections,
    )
except ImportError:  # pragma: no cover - supports direct script execution
    from process_pdfs import (
        build_document_metadata,
        discover_files,
        extract_docx_pages,
        extract_pdf_pages,
        extract_pptx_pages,
        normalize_extensions,
        normalize_text,
        select_files,
        strip_head_boilerplate,
        strip_vcsc_disclaimers,
        trim_tail_sections,
    )

LOGGER = logging.getLogger(__name__)
FIGURE_CAPTION_RE = re.compile(
    r"(?im)^\s*(figure|fig\.?|chart|exhibit|table|bảng|hình)\s*\d+[\.:\-]?\s+(.+)$"
)

# Lines matching any of these (after whitespace-normalization) are dropped entirely.
# Targets analyst-certification, compliance, ratings legends, contact/footer noise,
# bare numeric rows and other non-analytical chrome that hurts SFT/RAG quality.
LINE_NOISE_PATTERNS = [
    r"(?i)^see important disclosure.*$",
    r"(?i)^www\.vietcap\.com\.vn$",
    r"(?i)^page\s+\d+\s*(of\s+\d+)?$",
    r"(?i)^\d{1,4}\s*$",
    r"(?i)^analyst[s]?\s*certification.*$",
    r"(?i)^certification of analyst.*$",
    r"(?i)^disclaimer[s]?\s*$",
    r"(?i)^disclosures?\s*$",
    r"(?i)^important disclosures?.*$",
    r"(?i)^rating(?:\s+system|s)?\s*[:\-]?\s*$",
    r"(?i)^(buy|sell|hold|outperform|underperform|market\s*perform|neutral)\s*[:=].*$",
    r"(?i)^(mua|bán|nắm giữ|khả quan|kém khả quan|phù hợp thị trường)\s*[:=].*$",
    r"(?i)^(tel|fax|phone|email|e-mail|hotline)\s*[:.]?\s*\+?[\d\s().\-@a-z]+$",
    r"(?i)^\+?\(?\d[\d\s().\-]{6,}$",
    r"(?i)^[\w.\-]+@[\w.\-]+\.\w+$",
    r"(?i)^(head office|branch|chi nhánh|trụ sở|địa chỉ|address)\s*[:.].*$",
    r"(?i)^bloomberg\s*[:.].*$",
    r"(?i)^source\s*[:.]?\s*$",
    r"(?i)^nguồn\s*[:.]?\s*$",
    r"(?i)^[\d\s,.\-+%()$]+$",  # pure number / symbol rows orphaned from tables
    r"(?i)^(copyright|©|all rights reserved).*$",
    r"(?i)^vcsc\b.*\bresearch\b.*$",
]
_LINE_NOISE_RE = [re.compile(p) for p in LINE_NOISE_PATTERNS]

# A page is treated as pure boilerplate (dropped wholesale) when it contains
# several of these compliance/contact markers and little analytical text.
COMPLIANCE_PAGE_MARKERS = [
    "analyst certification",
    "analysts certification",
    "important disclosures",
    "rating system",
    "ratings definitions",
    "guarantee of future performance",
    "investment ratings",
    "this report is not directed",
    "khuyến cáo",
    "cam kết của chuyên viên",
    "hệ thống khuyến nghị",
    "miễn trách nhiệm",
]


def page_is_compliance(text: str, marker_threshold: int = 2) -> bool:
    lowered = text.casefold()
    hits = sum(1 for marker in COMPLIANCE_PAGE_MARKERS if marker in lowered)
    return hits >= marker_threshold


def content_fingerprint(markdown: str) -> str:
    """Hash the analytical body (excludes metadata header) for near-dup detection."""
    body = markdown.split("## Report Body", 1)[-1]
    collapsed = re.sub(r"\s+", " ", body).strip().casefold()
    return hashlib.sha1(collapsed.encode("utf-8")).hexdigest()


def parse_args() -> argparse.Namespace:
    repo_root = Path(__file__).resolve().parents[1]
    parser = argparse.ArgumentParser(
        description=(
            "Export cleaned report Markdown for PageIndex/vectorless RAG. "
            "The output preserves document/page structure instead of artificial chunks."
        )
    )
    parser.add_argument(
        "--input-dir",
        type=Path,
        default=repo_root / "raw_dataset",
        help="Root directory containing source reports (default: raw_dataset).",
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=repo_root / "ocr_pipeline" / "markdown_reports",
        help="Directory for cleaned Markdown files and manifest.jsonl.",
    )
    parser.add_argument(
        "--extensions",
        nargs="+",
        default=[".pdf", ".docx", ".pptx"],
        help="File extensions to include (default: .pdf .docx .pptx).",
    )
    parser.add_argument("--limit", type=int, default=None, help="Pilot file limit.")
    parser.add_argument(
        "--sample-mode",
        choices=["head", "random"],
        default="random",
        help="Selection mode when --limit is set (default: random).",
    )
    parser.add_argument("--seed", type=int, default=3407)
    parser.add_argument(
        "--trim-tail-pages",
        type=int,
        default=0,
        help="Drop this many trailing pages/slides/sections after boilerplate detection.",
    )
    parser.add_argument(
        "--min-page-words",
        type=int,
        default=20,
        help="Drop tiny pages after cleanup unless they look like headings.",
    )
    parser.add_argument(
        "--min-doc-words",
        type=int,
        default=80,
        help=(
            "Drop whole documents whose analytical body has fewer than this many words "
            "after cleanup. Removes near-empty caption-only shells (default: 80)."
        ),
    )
    parser.add_argument(
        "--keep-duplicates",
        action="store_true",
        help=(
            "Keep documents with identical analytical bodies. By default near-identical "
            "duplicates (e.g. PDF+PPTX+revised copies of the same report) are dropped."
        ),
    )
    parser.add_argument(
        "--use-markitdown",
        action="store_true",
        help=(
            "Use Microsoft's markitdown package when installed. It can improve layout "
            "for some Office/PDF files, but scanned PDFs still need OCR first."
        ),
    )
    parser.add_argument(
        "--log-level",
        default="INFO",
        choices=["DEBUG", "INFO", "WARNING", "ERROR"],
    )
    return parser.parse_args()


def slugify(value: str, fallback: str = "report") -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", value).strip("-").lower()
    return slug[:160] or fallback


def markdown_escape(text: str) -> str:
    return text.replace("\\", "\\\\").replace("`", "\\`")


def looks_like_heading(text: str) -> bool:
    if "|" in text:
        return False
    words = text.split()
    if not 1 <= len(words) <= 14:
        return False
    if text.endswith(('.', ',', ';', ':')):
        return False
    alpha = [c for c in text if c.isalpha()]
    if not alpha:
        return False
    upper_ratio = sum(c.isupper() for c in alpha) / len(alpha)
    return upper_ratio >= 0.45 or text.istitle() or text[0].isupper()


def clean_line(line: str) -> str:
    line = re.sub(r"\s+", " ", line).strip()
    if not line:
        return ""
    for pattern_re in _LINE_NOISE_RE:
        if pattern_re.match(line):
            return ""
    return line.strip()


def page_to_markdown(page_text: str) -> str:
    lines = [clean_line(line) for line in page_text.splitlines()]
    lines = [line for line in lines if line]
    blocks: list[str] = []
    paragraph: list[str] = []

    for line in lines:
        if looks_like_heading(line):
            if paragraph:
                blocks.append(" ".join(paragraph))
                paragraph = []
            blocks.append(f"### {line}")
        else:
            paragraph.append(line)

    if paragraph:
        blocks.append(" ".join(paragraph))

    return "\n\n".join(blocks).strip()


def markdown_table(rows: list[list[str]]) -> str:
    cleaned_rows = [[normalize_text(cell) for cell in row] for row in rows]
    cleaned_rows = [row for row in cleaned_rows if any(row)]
    if not cleaned_rows:
        return ""
    width = max(len(row) for row in cleaned_rows)
    padded = [row + [""] * (width - len(row)) for row in cleaned_rows]
    header = padded[0]
    body = padded[1:]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def extract_caption_notes(text: str) -> list[str]:
    notes: list[str] = []
    for match in FIGURE_CAPTION_RE.finditer(text):
        label = normalize_text(match.group(1)).title()
        caption = normalize_text(match.group(2))
        if caption:
            notes.append(f"{label}: {caption}")
    return notes


def infer_visual_kind(caption: str) -> str:
    lowered = caption.casefold()
    if any(term in lowered for term in ["breakdown", "structure", "market share", "mix"]):
        return "composition chart or breakdown figure"
    if any(term in lowered for term in ["trend", "growth", "cagr", "over time", "since"]):
        return "trend chart"
    if any(term in lowered for term in ["forecast", "estimate", "f", "projection"]):
        return "forecast chart/table"
    if any(term in lowered for term in ["correlation", "vs", "versus"]):
        return "comparison chart"
    if any(term in lowered for term in ["capacity", "volume", "production"]):
        return "operating metric chart"
    return "figure/chart"


def visual_notes_from_captions(text: str) -> str:
    notes = []
    for note in extract_caption_notes(text):
        caption = note.split(": ", 1)[-1]
        notes.append(f"- {note} ({infer_visual_kind(caption)})")
    if not notes:
        return ""
    return "### Figure and Chart Notes\n\n" + "\n".join(notes)


def format_page_parts(text: str, tables: list[list[list[str]]] | None = None, visual_note: str = "") -> str:
    parts = [page_to_markdown(text)] if text.strip() else []
    if visual_note.strip():
        parts.append(visual_note.strip())
    for table_index, table_rows in enumerate(tables or [], start=1):
        table_md = markdown_table(table_rows)
        if table_md:
            parts.append(f"### Extracted Table {table_index}\n\n{table_md}")
    return "\n\n".join(part for part in parts if part).strip()


def pages_to_markdown(
    pages: list[str],
    title: str,
    relative_source: str,
    min_page_words: int,
) -> tuple[str, int]:
    parts = [f"# {title}", "", "## Document Metadata", ""]
    parts.extend(
        [
            f"- Source: `{markdown_escape(relative_source)}`",
            "- Converted for: PageIndex vectorless RAG / SFT training",
            "- Cleanup: contact/compliance pages, analyst certification, ratings legends, "
            "disclosures, footers, and bare numeric rows removed where detected",
            "",
            "## Report Body",
        ]
    )

    retained_pages = 0
    body_words = 0
    for page_number, page_text in enumerate(pages, start=1):
        cleaned = strip_head_boilerplate(page_text)
        cleaned = strip_vcsc_disclaimers(cleaned)
        if page_is_compliance(cleaned):
            continue
        normalized = normalize_text(cleaned)
        if not normalized:
            continue
        if len(normalized.split()) < min_page_words and not looks_like_heading(normalized):
            continue

        if "\n### " in cleaned or "| ---" in cleaned:
            page_md = cleaned.strip()
        else:
            page_md = page_to_markdown(cleaned)
        if not page_md:
            continue
        retained_pages += 1
        body_words += len(normalize_text(page_md).split())
        parts.extend(["", f"## Page {page_number}", "", page_md])

    if retained_pages == 0:
        return "", 0
    return "\n".join(parts).strip() + "\n", body_words


def extract_with_markitdown(file_path: Path) -> list[str] | None:
    try:
        from markitdown import MarkItDown
    except ImportError:
        return None

    markdown = MarkItDown().convert(str(file_path)).text_content
    if not markdown.strip():
        return None
    return [markdown]


def extract_pdf_rich_pages(file_path: Path) -> list[str]:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("Missing dependency 'PyMuPDF'. Install with: pip install PyMuPDF") from exc

    pages: list[str] = []
    with fitz.open(file_path) as pdf_doc:
        for page in pdf_doc:
            text = page.get_text("text")
            tables: list[list[list[str]]] = []
            finder = getattr(page, "find_tables", None)
            if callable(finder):
                try:
                    for table in finder().tables:
                        rows = table.extract()
                        if rows:
                            tables.append([[str(cell or "") for cell in row] for row in rows])
                except Exception as exc:  # noqa: PERF203
                    LOGGER.debug("PDF table extraction skipped for %s: %s", file_path, exc)

            image_count = len(page.get_images(full=True))
            drawing_count = len(page.get_drawings())
            visual_note = visual_notes_from_captions(text)
            if image_count or drawing_count:
                summary = (
                    f"### Page Visual Summary\n\n"
                    f"- Detected embedded visuals: {image_count} image(s), {drawing_count} vector drawing object(s). "
                    "No image files were exported; use nearby captions and extracted tables as text-only figure context."
                )
                visual_note = "\n\n".join(part for part in [visual_note, summary] if part)

            pages.append(format_page_parts(text, tables=tables, visual_note=visual_note))
    return pages


def extract_docx_rich_pages(file_path: Path) -> list[str]:
    try:
        import docx
    except ImportError as exc:
        raise RuntimeError("Missing dependency 'python-docx'. Install with: pip install python-docx") from exc

    doc = docx.Document(file_path)
    blocks: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            blocks.append(text)
    for table_index, table in enumerate(doc.tables, start=1):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        table_md = markdown_table(rows)
        if table_md:
            blocks.append(f"Extracted Table {table_index}\n{table_md}")

    pages: list[str] = []
    for i in range(0, len(blocks), 30):
        text = "\n".join(blocks[i : i + 30])
        pages.append(format_page_parts(text, visual_note=visual_notes_from_captions(text)))
    return pages


def extract_pptx_rich_pages(file_path: Path) -> list[str]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Missing dependency 'python-pptx'. Install with: pip install python-pptx") from exc

    presentation = Presentation(file_path)
    slides: list[str] = []
    for slide in presentation.slides:
        text_parts: list[str] = []
        tables: list[list[list[str]]] = []
        chart_count = 0
        picture_count = 0
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                chart_count += 1
            if getattr(shape, "shape_type", None) is not None and "PICTURE" in str(getattr(shape, "shape_type", "")):
                picture_count += 1
            if getattr(shape, "has_table", False):
                table: Any = shape.table
                tables.append([[cell.text for cell in row.cells] for row in table.rows])
            text = getattr(shape, "text", "")
            if text and text.strip():
                text_parts.append(text.strip())

        text = "\n".join(text_parts)
        visual_note = visual_notes_from_captions(text)
        if chart_count or picture_count:
            summary = (
                "### Slide Visual Summary\n\n"
                f"- Detected visuals: {chart_count} chart object(s), {picture_count} picture object(s). "
                "No image files were exported; this is a text-only visual inventory."
            )
            visual_note = "\n\n".join(part for part in [visual_note, summary] if part)
        slides.append(format_page_parts(text, tables=tables, visual_note=visual_note))
    return slides


def extract_report_pages(file_path: Path, use_markitdown: bool) -> list[str]:
    if use_markitdown:
        pages = extract_with_markitdown(file_path)
        if pages:
            return pages
        LOGGER.warning("markitdown unavailable/empty; falling back to built-in parser: %s", file_path)

    extension = file_path.suffix.lower()
    if extension == ".pdf":
        pages = extract_pdf_rich_pages(file_path)
    elif extension == ".docx":
        pages = extract_docx_rich_pages(file_path)
    elif extension == ".pptx":
        pages = extract_pptx_rich_pages(file_path)
    else:
        raise ValueError(f"Unsupported extension: {extension}")
    return [page for page in pages if page.strip()]


def export_markdown_reports(args: argparse.Namespace) -> None:
    if not args.input_dir.exists():
        raise FileNotFoundError(f"Input directory does not exist: {args.input_dir}")

    extensions = normalize_extensions(args.extensions)
    files = discover_files(args.input_dir, extensions)
    selected_files = select_files(files, args.limit, args.sample_mode, args.seed)
    if not selected_files:
        raise RuntimeError("No files matched the requested input/extensions")

    args.output_dir.mkdir(parents=True, exist_ok=True)
    manifest_path = args.output_dir / "manifest.jsonl"
    failure_path = args.output_dir / "markdown_failures.log"
    manifest_rows: list[dict] = []
    failures: list[str] = []
    seen_fingerprints: dict[str, str] = {}
    dropped_short = 0
    dropped_duplicate = 0
    dropped_empty = 0

    for index, file_path in enumerate(selected_files, start=1):
        try:
            pages = extract_report_pages(file_path, args.use_markitdown)
            pages = trim_tail_sections(pages, args.trim_tail_pages)
            full_text = strip_vcsc_disclaimers("\n\n".join(pages))
            metadata = build_document_metadata(file_path, args.input_dir, full_text)
            markdown, body_words = pages_to_markdown(
                pages=pages,
                title=metadata["title"],
                relative_source=metadata["relative_source"],
                min_page_words=args.min_page_words,
            )
            if not markdown:
                dropped_empty += 1
                LOGGER.info("[%d/%d] Dropped (empty after cleanup): %s", index, len(selected_files), file_path.name)
                continue
            if body_words < args.min_doc_words:
                dropped_short += 1
                LOGGER.info(
                    "[%d/%d] Dropped (body %d < min_doc_words %d): %s",
                    index, len(selected_files), body_words, args.min_doc_words, file_path.name,
                )
                continue

            fingerprint = content_fingerprint(markdown)
            if not args.keep_duplicates and fingerprint in seen_fingerprints:
                dropped_duplicate += 1
                LOGGER.info(
                    "[%d/%d] Dropped (duplicate of %s): %s",
                    index, len(selected_files), seen_fingerprints[fingerprint], file_path.name,
                )
                continue

            output_name = f"{slugify(metadata['doc_id'])}.md"
            output_path = args.output_dir / output_name
            output_path.write_text(markdown, encoding="utf-8")
            seen_fingerprints[fingerprint] = output_name
            manifest_rows.append(
                {
                    **metadata,
                    "markdown_path": str(output_path.relative_to(args.output_dir)),
                    "markdown_word_count": len(markdown.split()),
                    "body_word_count": body_words,
                    "content_fingerprint": fingerprint,
                }
            )
            LOGGER.info("[%d/%d] Wrote %s", index, len(selected_files), output_path.name)
        except Exception as exc:  # noqa: PERF203
            failures.append(str(file_path))
            LOGGER.error("[%d/%d] Failed %s: %s", index, len(selected_files), file_path, exc)

    with manifest_path.open("w", encoding="utf-8") as manifest_file:
        for row in manifest_rows:
            manifest_file.write(json.dumps(row, ensure_ascii=False) + "\n")

    if failures:
        failure_path.write_text("\n".join(failures), encoding="utf-8")
        LOGGER.warning("Failed file list written to %s", failure_path)

    LOGGER.info(
        "Markdown reports: %d | Dropped: empty=%d short=%d duplicate=%d | Failures: %d",
        len(manifest_rows), dropped_empty, dropped_short, dropped_duplicate, len(failures),
    )
    LOGGER.info("Manifest: %s", manifest_path)


def main() -> int:
    args = parse_args()
    logging.basicConfig(
        level=getattr(logging, args.log_level),
        format="%(asctime)s | %(levelname)s | %(message)s",
    )
    try:
        export_markdown_reports(args)
    except Exception as exc:  # noqa: PERF203
        LOGGER.error("Markdown export failed: %s", exc)
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
