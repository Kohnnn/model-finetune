from __future__ import annotations

import argparse
import hashlib
import json
import logging
import random
import re
import sys
import traceback
from pathlib import Path
from time import perf_counter
from typing import Any, Iterable

LOGGER = logging.getLogger(__name__)

VIETNAMESE_CHAR_PATTERN = re.compile(
    r"[ăâđêôơưáàảãạấầẩẫậắằẳẵặéèẻẽẹếềểễệíìỉĩịóòỏõọốồổỗộớờởỡợúùủũụứừửữựýỳỷỹỵ]",
    re.IGNORECASE,
)
YEAR_PATTERN = re.compile(r"(?:19|20)\d{2}")

TAIL_SECTION_MARKERS = [
    "analyst certification",
    "analyst certification of independence",
    "important disclosures",
    "vCSC rating system",
    "valuation methodology",
    "for investment advice, trade execution or other enquiries",
    "for investment advice",
    "this report is provided, for information purposes only",
    "major us institutional investors",
    "u.k. and european economic area",
    "hong kong:",
    "new zealand:",
    "contacts",
    "contact us",
    "xác nhận của chuyên viên phân tích",
    "phương pháp định giá và hệ thống khuyến nghị của vcsc",
    "phương pháp định giá: để xác định giá mục tiêu",
    "phòng giao dịch chứng khoán",
    "phòng nghiên cứu và phân tích",
    "báo cáo này được viết và phát hành bởi",
    "công ty cổ phần chứng khoán bản việt",
]

CONTACT_PAGE_MARKERS = [
    "contacts",
    "contact us",
    "local sales representative",
    "for investment advice, trade execution or other enquiries",
    "decker&co",
]

HEAD_SECTION_MARKERS = [
    "update report",
    "see important disclosure",
    "www.vietcap.com.vn",
    "vietcap securities",
    "disclaimer",
]

NOISE_PATTERNS = [
    r"(?i)^\s*figure\s*\d+.*?source:\s*",
    r"(?i)^\s*bảng\s*\d+.*?nguồn:\s*",
    r"(?i)^\s*table\s*\d+.*?source:\s*",
    r"(?i)^\s*page\s*\d+\s+of\s+\d+",
    r"(?i)^\s*\d{1,3}\.\s+[a-z\-\s]+$",
    r"(?i)^\s*mục\s*lục\s*c",
    r"(?i)^\s*tài\s*liệu\s*tham\s*khảo",
    r"(?i)nguồn:\s*\w+",
    r"(?i)ngừng\s*theo\s*dõi",
    r"(?i)báo\s*cáo\s*tài\s*chính\s*nguồn:",
    r"(?i)^\s*[+\-]?\$?[\d,\.]+\s*%?$",
    r"(?i)^\s*[+\-]?\d+\.?\d*%\s*$",
    r"(?i)(buy|sell|outperform|underperform|neutral)\s*=\s*",
    r"(?i)vietcap rating system",
    r"(?i)^\+?\d{6,}",
    r"(?i)^\s*[ 	]+$",
]

_NOISE_PATTERN_RE = [re.compile(p) for p in NOISE_PATTERNS]

ANALYTICAL_MARKERS = [
    "target price",
    "valuation",
    "recommendation",
    "earnings",
    "margin",
    "profit",
    "forecast",
    "upside",
    "downside",
    "khuyến nghị",
    "giá mục tiêu",
    "định giá",
    "lợi nhuận",
    "biên lợi nhuận",
    "dự báo",
]

SYSTEM_PROMPT = (
    "You are a senior equity research analyst. Answer queries using a highly "
    "professional financial research tone, voice, and style, based strictly "
    "on the provided context. Focus on analytical synthesis and strategic "
    "insights rather than listing facts."
)

USER_TASK_PROMPT = (
    "Task: Deliver expert equity research commentary and strategic evaluation "
    "based on the context above. Prioritize deep analysis over factual reporting."
)


def parse_args() -> argparse.Namespace:
    repo_root = Path(__file__).resolve().parents[1]

    parser = argparse.ArgumentParser(
        description=(
            "Parse dataset documents into RAG chunks and finetuning templates "
            "(JSONL output)."
        )
    )
    parser.add_argument(
        "--input-dir",
        type=Path,
        default=repo_root / "raw_dataset",
        help="Root directory containing source files (default: raw_dataset).",
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=Path(__file__).resolve().parent,
        help=(
            "Directory for chroma_chunks.jsonl and finetune_template.jsonl "
            "(default: ocr_pipeline)."
        ),
    )
    parser.add_argument(
        "--extensions",
        nargs="+",
        default=[".pdf", ".docx", ".pptx"],
        help=(
            "File extensions to include, e.g. .pdf .docx .pptx "
            "(default: .pdf .docx .pptx)."
        ),
    )
    parser.add_argument(
        "--limit",
        type=int,
        default=None,
        help="Limit number of files processed (useful for pilot runs).",
    )
    parser.add_argument(
        "--sample-mode",
        choices=["head", "random"],
        default="random",
        help="How files are selected when --limit is set (default: random).",
    )
    parser.add_argument(
        "--seed",
        type=int,
        default=3407,
        help="Random seed for deterministic pilot sampling (default: 3407).",
    )
    parser.add_argument(
        "--chunk-words",
        type=int,
        default=800,
        help="Chunk size in words (default: 800).",
    )
    parser.add_argument(
        "--overlap-words",
        type=int,
        default=100,
        help="Word overlap between chunks (default: 100).",
    )
    parser.add_argument(
        "--min-chunk-words",
        type=int,
        default=200,
        help="Discard chunks smaller than this threshold (default: 200).",
    )
    parser.add_argument(
        "--trim-tail-pages",
        type=int,
        default=0,
        help=(
            "Drop this many trailing pages/slides/sections as likely disclaimers "
            "when possible (default: 0)."
        ),
    )
    parser.add_argument(
        "--keep-duplicates",
        action="store_true",
        help="Keep documents with identical normalized analytical content.",
    )
    parser.add_argument(
        "--near-duplicate-distance",
        type=int,
        default=3,
        help="Maximum 64-bit SimHash distance for grouping revisions into one document family.",
    )
    parser.add_argument(
        "--allow-pilot-overwrite",
        action="store_true",
        help="Allow --limit to overwrite the default ocr_pipeline outputs.",
    )
    parser.add_argument(
        "--log-level",
        default="INFO",
        choices=["DEBUG", "INFO", "WARNING", "ERROR"],
        help="Logging verbosity (default: INFO).",
    )
    return parser.parse_args()


def normalize_extensions(values: list[str]) -> list[str]:
    normalized: list[str] = []
    for value in values:
        ext = value.strip().lower()
        if not ext.startswith("."):
            ext = f".{ext}"
        normalized.append(ext)
    return sorted(set(normalized))


def discover_files(input_dir: Path, extensions: list[str]) -> list[Path]:
    files: list[Path] = []
    for extension in extensions:
        files.extend(input_dir.rglob(f"*{extension}"))
    return sorted(
        (p for p in files if p.is_file() and not should_skip_file(p)),
        key=lambda p: str(p).lower(),
    )


def should_skip_file(file_path: Path) -> bool:
    return file_path.name.startswith("~$")


def select_files(
    files: list[Path],
    limit: int | None,
    sample_mode: str,
    seed: int,
) -> list[Path]:
    if limit is None or limit >= len(files):
        return files
    if sample_mode == "head":
        return files[:limit]
    rng = random.Random(seed)
    return sorted(rng.sample(files, limit), key=lambda p: str(p).lower())


def normalize_text(text: str) -> str:
    compact = re.sub(r"\s+", " ", text)
    return compact.strip()


def normalize_for_matching(text: str) -> str:
    return normalize_text(text).casefold()


def count_marker_hits(text: str, markers: list[str]) -> int:
    normalized = normalize_for_matching(text)
    return sum(normalize_for_matching(marker) in normalized for marker in markers)


def strip_head_boilerplate(text: str) -> str:
    text_lower = normalize_for_matching(text)
    found_marker = None
    found_idx = len(text)
    for marker in HEAD_SECTION_MARKERS:
        norm_marker = normalize_for_matching(marker)
        idx = text_lower.find(norm_marker)
        if idx != -1 and idx < found_idx:
            found_idx = idx
            found_marker = marker
    if found_marker is None:
        return text
    lines = text.split("\n")
    for i, line in enumerate(lines):
        line_lower = normalize_for_matching(line)
        if normalize_for_matching(found_marker) in line_lower:
            return "\n".join(lines[i + 1 :]).strip()
    return text


def matches_noise_pattern(text: str) -> bool:
    for pattern_re in _NOISE_PATTERN_RE:
        if pattern_re.search(text):
            return True
    return False


def has_excessive_numbers(text: str) -> bool:
    if not text:
        return False
    digits = sum(c.isdigit() for c in text)
    return (digits / len(text)) > 0.30


def count_analytical_markers(text: str) -> int:
    return count_marker_hits(text, ANALYTICAL_MARKERS)


def is_boilerplate_page(text: str) -> bool:
    word_count = len(text.split())
    tail_hits = count_marker_hits(text, TAIL_SECTION_MARKERS)
    contact_hits = count_marker_hits(text, CONTACT_PAGE_MARKERS)
    analytical_hits = count_marker_hits(text, ANALYTICAL_MARKERS)
    normalized = normalize_for_matching(text)

    if not normalized:
        return False
    if contact_hits >= 1 and word_count <= 220:
        return True
    if (
        "analyst certification" in normalized
        or "xác nhận của chuyên viên phân tích" in normalized
    ):
        return True
    if tail_hits >= 2 and analytical_hits == 0:
        return True
    if tail_hits >= 3:
        return True
    return False


def strip_vcsc_disclaimers(text: str) -> str:
    """Truncates disclaimer and contact boilerplate typically found at report tails."""

    text_lower = normalize_for_matching(text)
    earliest_idx = len(text)

    for marker in TAIL_SECTION_MARKERS:
        norm_marker = normalize_for_matching(marker)
        idx = text_lower.find(norm_marker)
        if idx != -1 and idx < earliest_idx:
            earliest_idx = idx

    if earliest_idx < len(text):
        return text[:earliest_idx].strip()
    return text


def infer_document_title(file_path: Path) -> str:
    title = re.sub(r"[_\-]+", " ", file_path.stem)
    title = re.sub(r"\s+", " ", title)
    return title.strip()


def infer_document_year(relative_source: str) -> int | None:
    matches = YEAR_PATTERN.findall(relative_source)
    if not matches:
        return None
    return int(matches[0])


def infer_document_language(relative_source: str, text: str) -> str:
    normalized_source = relative_source.lower()
    if "[vn]" in normalized_source or "vietnamese" in normalized_source:
        return "vi"
    if VIETNAMESE_CHAR_PATTERN.search(text):
        return "vi"
    return "en"


def build_doc_id(relative_source: str) -> str:
    normalized_source = relative_source.lower()
    doc_id = re.sub(r"[^a-z0-9]+", "_", normalized_source).strip("_")
    digest = hashlib.md5(normalized_source.encode("utf-8")).hexdigest()[:8]
    return f"{doc_id}_{digest}"


def sha256_file(file_path: Path) -> str:
    digest = hashlib.sha256()
    with file_path.open("rb") as source_file:
        for block in iter(lambda: source_file.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def content_sha256(text: str) -> str:
    return hashlib.sha256(normalize_for_matching(text).encode("utf-8")).hexdigest()


def content_simhash(text: str) -> int:
    tokens = re.findall(r"\w+", normalize_for_matching(text))
    if not tokens:
        return 0
    shingles = [" ".join(tokens[index : index + 5]) for index in range(max(1, len(tokens) - 4))]
    weights = [0] * 64
    for shingle in shingles:
        value = int.from_bytes(hashlib.blake2b(shingle.encode("utf-8"), digest_size=8).digest(), "big")
        for bit in range(64):
            weights[bit] += 1 if value & (1 << bit) else -1
    result = 0
    for bit, weight in enumerate(weights):
        if weight >= 0:
            result |= 1 << bit
    return result


def simhash_distance(first: int, second: int) -> int:
    return (first ^ second).bit_count()


def find_family_root(parents: dict[str, str], family_id: str) -> str:
    root = family_id
    while parents[root] != root:
        root = parents[root]
    while parents[family_id] != family_id:
        parent = parents[family_id]
        parents[family_id] = root
        family_id = parent
    return root


def add_near_duplicate_family(
    simhash: int,
    family_id: str,
    families: list[tuple[int, str]],
    parents: dict[str, str],
    maximum_distance: int,
) -> bool:
    parents.setdefault(family_id, family_id)
    matches = {
        candidate_family
        for candidate_simhash, candidate_family in families
        if simhash_distance(simhash, candidate_simhash) <= maximum_distance
    }
    for candidate_family in matches:
        first = find_family_root(parents, family_id)
        second = find_family_root(parents, candidate_family)
        if first != second:
            canonical, merged = sorted((first, second))
            parents[merged] = canonical
    families.append((simhash, family_id))
    return bool(matches)


def canonicalize_document_families(
    rows: list[dict[str, Any]],
    parents: dict[str, str],
) -> None:
    for row in rows:
        metadata = row.get("metadata") if isinstance(row.get("metadata"), dict) else row
        family_id = metadata.get("document_family_id")
        if isinstance(family_id, str) and family_id in parents:
            metadata["document_family_id"] = find_family_root(parents, family_id)


def build_document_metadata(
    file_path: Path,
    input_dir: Path,
    text: str,
    extraction_method: str = "primary",
    extraction_warnings: list[str] | None = None,
) -> dict[str, Any]:
    relative_source = str(file_path.relative_to(input_dir))
    normalized_sha256 = content_sha256(text)
    return {
        "relative_source": relative_source,
        "doc_id": build_doc_id(relative_source),
        "title": infer_document_title(file_path),
        "year": infer_document_year(relative_source),
        "language": infer_document_language(relative_source, text),
        "file_extension": file_path.suffix.lower(),
        "source_file_sha256": sha256_file(file_path),
        "content_sha256": normalized_sha256,
        "document_family_id": f"family_{normalized_sha256}",
        "extraction_method": extraction_method,
        "parser_schema_version": "2",
        "extraction_warnings": json.dumps(extraction_warnings or [], ensure_ascii=False),
    }


def extract_pdf_pages(file_path: Path) -> list[str]:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError(
            "Missing dependency 'PyMuPDF'. Install with: pip install PyMuPDF"
        ) from exc

    pages: list[str] = []
    with fitz.open(file_path) as pdf_doc:
        for page in pdf_doc:
            pages.append(page.get_text("text"))
    return pages


def extract_docx_pages(file_path: Path) -> list[str]:
    try:
        import docx
    except ImportError as exc:
        raise RuntimeError(
            "Missing dependency 'python-docx'. Install with: pip install python-docx"
        ) from exc

    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = docx.Document(file_path)
    blocks: list[str] = []
    for child in doc.element.body.iterchildren():
        if child.tag.endswith("}p"):
            text = Paragraph(child, doc).text.strip()
            if text:
                blocks.append(text)
            continue
        if child.tag.endswith("}tbl"):
            table = Table(child, doc)
            rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
            rows = [row for row in rows if row.strip(" |")]
            if rows:
                blocks.append("Table:\n" + "\n".join(rows))
    return ["\n".join(blocks[i : i + 30]) for i in range(0, len(blocks), 30)]


def chart_text(chart: Any) -> list[str]:
    parts: list[str] = []
    if chart.has_title:
        title = chart.chart_title.text_frame
        if title and title.text.strip():
            parts.append(title.text.strip())
    for series in chart.series:
        values = [str(value) for value in series.values]
        categories: list[str] = []
        try:
            categories = [str(category.label) for category in chart.plots[0].categories]
        except (AttributeError, IndexError, TypeError):
            pass
        pairs = [
            f"{category}: {value}"
            for category, value in zip(categories, values)
        ]
        parts.append(f"Chart series {series.name}: " + "; ".join(pairs or values))
    return [part for part in parts if part.strip()]


def extract_pptx_pages(file_path: Path) -> list[str]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "Missing dependency 'python-pptx'. Install with: pip install python-pptx"
        ) from exc

    presentation = Presentation(file_path)
    slides: list[str] = []
    for slide in presentation.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows]
                parts.extend(row for row in rows if row.strip(" |"))
            else:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    parts.append(text.strip())
            if getattr(shape, "has_chart", False):
                parts.extend(chart_text(shape.chart))
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except (AttributeError, ValueError):
            notes = ""
        if notes:
            parts.append(f"Speaker notes: {notes}")
        slides.append("\n".join(parts))
    return slides


def extract_document_pages(file_path: Path) -> tuple[list[tuple[int, str]], list[str]]:
    extension = file_path.suffix.lower()
    warnings: list[str] = []
    if extension == ".pdf":
        pages = extract_pdf_pages(file_path)
        try:
            import fitz

            with fitz.open(file_path) as pdf_doc:
                for page_number, (page, text) in enumerate(zip(pdf_doc, pages), start=1):
                    if len(normalize_text(text).split()) < 5 and page.get_images(full=True):
                        warnings.append(f"likely_scanned_or_image_only_page:{page_number}")
        except ImportError:
            pass
    elif extension == ".docx":
        pages = extract_docx_pages(file_path)
    elif extension == ".pptx":
        pages = extract_pptx_pages(file_path)
    else:
        raise ValueError(f"Unsupported extension: {extension}")
    return [
        (number, normalized)
        for number, page in enumerate(pages, start=1)
        if (normalized := normalize_text(page))
    ], warnings


def extract_pages(file_path: Path) -> list[str]:
    pages, _ = extract_document_pages(file_path)
    return [text for _, text in pages]


def trim_tail_sections(pages: list[str], trim_tail_pages: int) -> list[str]:
    if not pages:
        return []

    trimmed_pages = list(pages)
    while len(trimmed_pages) > 1 and is_boilerplate_page(trimmed_pages[-1]):
        trimmed_pages.pop()

    if trim_tail_pages <= 0:
        return trimmed_pages
    if len(trimmed_pages) > trim_tail_pages:
        return trimmed_pages[:-trim_tail_pages]
    if len(trimmed_pages) > 1:
        return trimmed_pages[:-1]
    return trimmed_pages


def is_boilerplate_chunk(text: str) -> bool:
    normalized = normalize_for_matching(text)
    if not normalized:
        return True

    word_count = len(text.split())
    tail_hits = count_marker_hits(text, TAIL_SECTION_MARKERS)
    contact_hits = count_marker_hits(text, CONTACT_PAGE_MARKERS)
    analytical_hits = count_analytical_markers(text)

    if contact_hits >= 1 and word_count <= 220:
        return True
    if (
        "analyst certification" in normalized
        or "xác nhận của chuyên viên phân tích" in normalized
    ):
        return True
    if tail_hits >= 2 and analytical_hits == 0:
        return True
    if tail_hits >= 3:
        return True
    if matches_noise_pattern(text):
        return True
    if has_excessive_numbers(text):
        return True
    return False


def is_quality_chunk(text: str) -> bool:
    if is_boilerplate_chunk(text):
        return False
    word_count = len(text.split())
    if word_count < 200:
        return False
    analytical_count = count_analytical_markers(text)
    if analytical_count < 2:
        return False
    return True


def chunk_text(
    text: str,
    chunk_words: int,
    overlap_words: int,
    min_chunk_words: int,
) -> list[str]:
    if overlap_words >= chunk_words:
        raise ValueError("overlap_words must be smaller than chunk_words")

    words = text.split()
    if not words:
        return []

    chunks: list[str] = []
    step = chunk_words - overlap_words
    for start in range(0, len(words), step):
        chunk = words[start : start + chunk_words]
        if len(chunk) >= min_chunk_words:
            chunks.append(" ".join(chunk))
    return chunks


def split_sentences(text: str) -> list[list[str]]:
    sentences = re.split(r"(?<=[.!?])\s+|\n+", text)
    return [sentence.split() for sentence in sentences if sentence.split()]


def chunk_sections(
    pages: list[tuple[int, str]],
    chunk_words: int,
    overlap_words: int,
    min_chunk_words: int,
) -> list[dict[str, Any]]:
    if overlap_words >= chunk_words:
        raise ValueError("overlap_words must be smaller than chunk_words")
    units: list[tuple[int, list[str], int]] = []
    source_word = 0
    for page_number, page_text in pages:
        for sentence in split_sentences(page_text):
            while sentence:
                part, sentence = sentence[:chunk_words], sentence[chunk_words:]
                units.append((page_number, part, source_word))
                source_word += len(part)
    chunks: list[dict[str, Any]] = []
    current: list[tuple[int, list[str], int]] = []
    for unit in units:
        current_words = sum(len(words) for _, words, _ in current)
        if current and current_words + len(unit[1]) > chunk_words:
            words = [word for _, sentence, _ in current for word in sentence]
            if len(words) >= min_chunk_words:
                chunks.append(
                    {
                        "text": " ".join(words),
                        "start_page": current[0][0],
                        "end_page": current[-1][0],
                        "source_page_numbers": json.dumps(
                            list(dict.fromkeys(page for page, _, _ in current))
                        ),
                        "source_word_start": current[0][2],
                        "source_word_end": current[-1][2] + len(current[-1][1]),
                    }
                )
            overlap: list[tuple[int, list[str], int]] = []
            overlap_count = 0
            for previous in reversed(current):
                if overlap_count + len(previous[1]) > overlap_words:
                    break
                overlap.insert(0, previous)
                overlap_count += len(previous[1])
            current = overlap
        current.append(unit)
    if current:
        words = [word for _, sentence, _ in current for word in sentence]
        if len(words) >= min_chunk_words:
            chunks.append(
                {
                    "text": " ".join(words),
                    "start_page": current[0][0],
                    "end_page": current[-1][0],
                    "source_page_numbers": json.dumps(
                        list(dict.fromkeys(page for page, _, _ in current))
                    ),
                    "source_word_start": current[0][2],
                    "source_word_end": current[-1][2] + len(current[-1][1]),
                }
            )
    return chunks


def write_jsonl(path: Path, rows: Iterable[dict]) -> None:
    with path.open("w", encoding="utf-8") as output_file:
        for row in rows:
            output_file.write(json.dumps(row, ensure_ascii=False) + "\n")


def process_dataset(args: argparse.Namespace) -> None:
    start_time = perf_counter()
    default_output_dir = Path(__file__).resolve().parent
    if args.limit is not None and args.output_dir.resolve() == default_output_dir.resolve() and not args.allow_pilot_overwrite:
        raise RuntimeError("--limit requires --allow-pilot-overwrite when using the default output directory")
    if not args.input_dir.exists():
        raise FileNotFoundError(f"Input directory does not exist: {args.input_dir}")

    extensions = normalize_extensions(args.extensions)
    files = discover_files(args.input_dir, extensions)
    selected_files = select_files(files, args.limit, args.sample_mode, args.seed)
    if not selected_files:
        raise RuntimeError("No files matched the given input directory and extension filters")

    chroma_chunks: list[dict] = []
    finetune_templates: list[dict] = []
    manifest_rows: list[dict[str, Any]] = []
    failures: list[str] = []
    seen_content: dict[str, str] = {}
    family_simhashes: list[tuple[int, str]] = []
    family_parents: dict[str, str] = {}
    if not 0 <= args.near_duplicate_distance <= 64:
        raise ValueError("--near-duplicate-distance must be between 0 and 64")

    for index, file_path in enumerate(selected_files, start=1):
        file_start = perf_counter()
        relative_source = str(file_path.relative_to(args.input_dir))
        try:
            pages, warnings = extract_document_pages(file_path)
            trimmed_texts = trim_tail_sections([text for _, text in pages], args.trim_tail_pages)
            pages = pages[: len(trimmed_texts)]
            cleaned_pages = [
                (number, strip_vcsc_disclaimers(strip_head_boilerplate(text)))
                for number, text in pages
            ]
            cleaned_pages = [(number, text) for number, text in cleaned_pages if text]
            text = "\n\n".join(text for _, text in cleaned_pages)
            metadata = build_document_metadata(file_path, args.input_dir, text, extraction_warnings=warnings)
            content_hash = metadata["content_sha256"]
            if not args.keep_duplicates and content_hash in seen_content:
                manifest_rows.append({
                    **metadata,
                    "status": "skipped",
                    "reason_code": "duplicate_normalized_content",
                    "duplicate_of": seen_content[content_hash],
                    "chunk_count": 0,
                })
                continue
            seen_content[content_hash] = relative_source
            simhash = content_simhash(text)
            if add_near_duplicate_family(
                simhash,
                metadata["document_family_id"],
                family_simhashes,
                family_parents,
                args.near_duplicate_distance,
            ):
                warnings.append("near_duplicate_revision_grouped")
                metadata["extraction_warnings"] = json.dumps(warnings, ensure_ascii=False)
            chunks = chunk_sections(
                cleaned_pages, args.chunk_words, args.overlap_words, args.min_chunk_words
            )
            chunks = [chunk for chunk in chunks if is_quality_chunk(chunk["text"])]
            if not chunks:
                manifest_rows.append({
                    **metadata,
                    "status": "skipped",
                    "reason_code": "no_quality_chunks",
                    "chunk_count": 0,
                })
                continue
            for chunk_index, chunk in enumerate(chunks):
                chunk_id = f"{metadata['doc_id']}_chunk_{chunk_index:04d}"
                chunk_metadata = {
                    **metadata,
                    **{key: value for key, value in chunk.items() if key != "text"},
                    "chunk_index": chunk_index,
                    "chunk_word_count": len(chunk["text"].split()),
                }
                chroma_chunks.append({"id": chunk_id, "text": chunk["text"], "metadata": chunk_metadata})
                finetune_templates.append({
                    "metadata": {
                        **chunk_metadata,
                        "context_sha256": hashlib.sha256(chunk["text"].encode("utf-8")).hexdigest(),
                        "task_type": "analytical_synthesis",
                        "source_spans": [
                            {
                                "start_page": chunk["start_page"],
                                "end_page": chunk["end_page"],
                                "source_word_start": chunk["source_word_start"],
                                "source_word_end": chunk["source_word_end"],
                            }
                        ],
                        "review_status": "draft",
                        "reviewed_by": "",
                        "reviewed_at": "",
                        "approval_checklist_version": "v1",
                        "verified_external_numbers": [],
                    },
                    "messages": [
                        {"role": "system", "content": SYSTEM_PROMPT},
                        {"role": "user", "content": f"Context:\n{chunk['text']}\n\n{USER_TASK_PROMPT}"},
                        {"role": "assistant", "content": ""},
                    ],
                })
            manifest_rows.append({
                **metadata,
                "status": "parsed",
                "reason_code": "ok",
                "chunk_count": len(chunks),
                "elapsed_seconds": round(perf_counter() - file_start, 3),
            })
            LOGGER.info("[%d/%d] Parsed %s -> %d chunks", index, len(selected_files), file_path.name, len(chunks))
        except Exception as exc:  # noqa: PERF203
            details = traceback.format_exc()
            manifest_rows.append({
                "relative_source": relative_source,
                "status": "failed",
                "reason_code": type(exc).__name__,
                "exception": str(exc),
                "exception_details": details,
                "chunk_count": 0,
            })
            failures.append(f"{relative_source}\n{details}")
            LOGGER.error("[%d/%d] Failed %s: %s", index, len(selected_files), file_path.name, exc)

    canonicalize_document_families(chroma_chunks, family_parents)
    canonicalize_document_families(finetune_templates, family_parents)
    canonicalize_document_families(manifest_rows, family_parents)
    args.output_dir.mkdir(parents=True, exist_ok=True)
    write_jsonl(args.output_dir / "chroma_chunks.jsonl", chroma_chunks)
    write_jsonl(args.output_dir / "finetune_template.jsonl", finetune_templates)
    write_jsonl(args.output_dir / "parse_manifest.jsonl", manifest_rows)
    if failures:
        (args.output_dir / "parse_failures.log").write_text("\n\n".join(failures), encoding="utf-8")
    LOGGER.info("Run complete in %.2fs | chunks=%d", perf_counter() - start_time, len(chroma_chunks))


def main() -> int:
    args = parse_args()

    logging.basicConfig(
        level=getattr(logging, args.log_level),
        format="%(asctime)s | %(levelname)s | %(message)s",
    )

    try:
        process_dataset(args)
    except Exception as exc:  # noqa: PERF203
        LOGGER.error("Parsing failed: %s", exc)
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
